import uuid
from datetime import datetime
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from app.config import settings


def generate_professional_pptx(
    title: str,
    subtitle: str,
    slides_data: list[dict],
    file_prefix: str = "Executive_Briefing"
) -> str:
    """Generates an executive presentation deck (.pptx) with clean styling and layout."""
    doc_id = uuid.uuid4().hex[:6].upper()
    file_name = f"{file_prefix}_{datetime.utcnow().strftime('%Y%m%d')}_{doc_id}.pptx"
    file_path = settings.ARTIFACTS_DIR / file_name

    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9 Widescreen
    prs.slide_height = Inches(7.5)

    # 1. Title Slide
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)

    # Background accent bar
    top_bar = slide.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.333), Inches(3.5))
    tf = top_bar.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "MRPL SOVEREIGN INDUSTRIAL OPERATIONS"
    p.font.name = "Arial"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = RGBColor(30, 58, 138)

    p2 = tf.add_paragraph()
    p2.text = title
    p2.font.name = "Arial"
    p2.font.size = Pt(32)
    p2.font.bold = True
    p2.font.color.rgb = RGBColor(15, 23, 42)

    p3 = tf.add_paragraph()
    p3.text = f"{subtitle} | Generated: {datetime.utcnow().strftime('%d-%b-%Y')}"
    p3.font.name = "Arial"
    p3.font.size = Pt(16)
    p3.font.color.rgb = RGBColor(100, 116, 139)

    # 2. Content Slides
    for slide_info in slides_data:
        s = prs.slides.add_slide(blank_slide_layout)
        
        # Header
        hdr_box = s.shapes.add_textbox(Inches(1.0), Inches(0.8), Inches(11.333), Inches(1.0))
        htf = hdr_box.text_frame
        hp = htf.paragraphs[0]
        hp.text = slide_info.get("heading", "Executive Update")
        hp.font.name = "Arial"
        hp.font.size = Pt(24)
        hp.font.bold = True
        hp.font.color.rgb = RGBColor(30, 58, 138)

        # Body Bullets
        body_box = s.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.333), Inches(4.5))
        btf = body_box.text_frame
        btf.word_wrap = True

        bullets = slide_info.get("bullets", [])
        for idx, bullet in enumerate(bullets):
            bp = btf.paragraphs[0] if idx == 0 else btf.add_paragraph()
            bp.text = f"•  {bullet}"
            bp.font.name = "Arial"
            bp.font.size = Pt(18)
            bp.font.color.rgb = RGBColor(30, 41, 59)
            bp.space_after = Pt(14)

    prs.save(str(file_path))
    return str(file_path)
